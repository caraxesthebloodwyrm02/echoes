# core/memory_manager.py
from pathlib import Path
import hashlib
import mimetypes
import logging
from typing import Dict, Any, Optional, List, Union
import fitz  # PyMuPDF
from PIL import Image
import docx
from pptx import Presentation
import pandas as pd
from openpyxl import load_workbook

logger = logging.getLogger(__name__)


class FileProcessor:
    """Handles processing of different file types."""

    def __init__(self, cache_dir: Path):
        self.cache_dir = cache_dir
        self.cache_dir.mkdir(parents=True, exist_ok=True)

    def get_file_hash(self, file_path: Path) -> str:
        """Generate a hash of the file contents for caching."""
        hasher = hashlib.sha256()
        with open(file_path, "rb") as f:
            while chunk := f.read(8192):
                hasher.update(chunk)
        return hasher.hexdigest()

    def process_file(self, file_path: Path) -> Dict[str, Any]:
        """Process a file and return its content and metadata."""
        logger.info(f"Starting to process file: {file_path}")
        if not file_path.exists():
            error_msg = f"File not found: {file_path}"
            logger.error(error_msg)
            return {"error": error_msg}

        try:
            mime_type, _ = mimetypes.guess_type(file_path)
            logger.debug(f"Detected MIME type: {mime_type} for {file_path}")

            file_hash = self.get_file_hash(file_path)
            cache_file = self.cache_dir / f"{file_hash}.json"

            # Check cache first
            if cache_file.exists():
                logger.info(f"Loading from cache: {cache_file}")
                import json

                with open(cache_file, "r") as f:
                    return json.load(f)

            logger.info(f"Processing file: {file_path} (size: {file_path.stat().st_size / (1024*1024):.2f}MB)")
            if file_path.suffix.lower() == ".html":
                logger.info("HTML file detected, will extract text content")

            # Process based on file type
            if mime_type:
                if mime_type.startswith("text/"):
                    result = self._process_text_file(file_path)
                elif mime_type.startswith("image/"):
                    result = self._process_image_file(file_path)
                elif mime_type == "application/pdf":
                    result = self._process_pdf_file(file_path)
                elif mime_type in [
                    "application/msword",
                    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                ]:
                    result = self._process_word_file(file_path)
                elif mime_type in [
                    "application/vnd.ms-excel",
                    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                ]:
                    result = self._process_excel_file(file_path)
                elif mime_type in [
                    "application/vnd.ms-powerpoint",
                    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                ]:
                    result = self._process_powerpoint_file(file_path)
                else:
                    result = {"error": f"Unsupported file type: {mime_type}"}
            else:
                result = {"error": f"Could not determine file type: {file_path}"}

            # Add common metadata
            if "error" not in result:
                result.update(
                    {
                        "file_path": str(file_path),
                        "file_name": file_path.name,
                        "file_size": file_path.stat().st_size,
                        "mime_type": mime_type,
                        "last_modified": file_path.stat().st_mtime,
                    }
                )

                # Save to cache
                with open(cache_file, "w") as f:
                    import json

                    json.dump(result, f)

            return result

        except Exception as e:
            logger.error(f"Error processing file {file_path}: {str(e)}")
            return {"error": str(e)}

    def _process_text_file(self, file_path: Path, max_size_mb: int = 5) -> Dict[str, Any]:
        """Process a text file with memory-efficient handling for large files.

        Args:
            file_path: Path to the text file
            max_size_mb: Maximum file size in MB to process in memory (default: 5MB)

        Returns:
            Dictionary with content and metadata
        """
        logger.info(f"Processing text file: {file_path}")
        try:
            file_size_mb = file_path.stat().st_size / (1024 * 1024)
            logger.info(f"File size: {file_size_mb:.2f}MB")

            # For HTML files, use specialized HTML processing
            if file_path.suffix.lower() == ".html":
                # For HTML files larger than 1MB, always use the streaming parser
                # This is more memory efficient even for moderately sized HTML files
                if file_size_mb > 1.0:  # 1MB threshold for HTML files
                    logger.info(f"HTML file detected ({file_size_mb:.2f}MB), using streaming parser")
                    return self._process_large_html_file(file_path)

                # For small HTML files, use regular BeautifulSoup
                logger.info("Using BeautifulSoup to parse small HTML content")
                try:
                    from bs4 import BeautifulSoup

                    with open(file_path, "r", encoding="utf-8") as f:
                        soup = BeautifulSoup(f, "html.parser")
                        # Remove script and style elements
                        for script in soup(["script", "style"]):
                            script.decompose()
                        content = soup.get_text(separator="\n", strip=True)
                        logger.info(f"Extracted {len(content)} characters of text from HTML")
                        return {"content": content, "type": "text", "source": "html"}
                except Exception as e:
                    logger.warning(
                        f"Error processing HTML with BeautifulSoup: {str(e)}, falling back to text processing"
                    )
                    # Fall through to regular text processing

            # For regular text files
            if file_size_mb > max_size_mb:
                logger.warning(f"Large text file detected ({file_size_mb:.2f}MB), using streaming reader")
                return self._process_large_text_file(file_path)

            with open(file_path, "r", encoding="utf-8") as f:
                content = f.read()
                logger.info(f"Read {len(content)} characters from text file")
                return {"content": content, "type": "text"}

        except Exception as e:
            logger.error(f"Error processing text file {file_path}: {str(e)}", exc_info=True)
            return {"error": f"Error processing text file: {str(e)}"}

    def _process_large_html_file(self, file_path: Path, chunk_size: int = 64 * 1024) -> Dict[str, Any]:
        """Process large HTML files in a memory-efficient way using a streaming parser.

        Args:
            file_path: Path to the HTML file
            chunk_size: Size of chunks to read from file (in bytes)

        Returns:
            Dictionary with extracted content and metadata
        """
        try:
            from bs4 import BeautifulSoup
            from bs4.builder import HTMLParserTreeBuilder

            logger.info(f"Processing large HTML file with streaming parser: {file_path}")

            # Initialize a list to hold text chunks
            text_chunks = []

            # Create a custom builder that processes text in chunks
            class StreamingHTMLParser(HTMLParserTreeBuilder):
                def __init__(self, *args, **kwargs):
                    super().__init__(*args, **kwargs)
                    self.text_chunks = []
                    self.in_script_or_style = False

                def handle_starttag(self, name, attrs, **kwargs):
                    super().handle_starttag(name, attrs, **kwargs)
                    self.in_script_or_style = name.lower() in ("script", "style")

                def handle_endtag(self, name, **kwargs):
                    super().handle_endtag(name, **kwargs)
                    self.in_script_or_style = False

                def handle_data(self, data):
                    if not self.in_script_or_style and data.strip():
                        self.text_chunks.append(data.strip())

            # Process the file in chunks
            with open(file_path, "rb") as f:
                # Read the file in chunks and process with BeautifulSoup
                parser = "html.parser"
                soup = BeautifulSoup("", parser, builder=StreamingHTMLParser())

                # Process file in chunks
                while True:
                    chunk = f.read(chunk_size)
                    if not chunk:
                        break

                    try:
                        # Try UTF-8 first
                        chunk_str = chunk.decode("utf-8")
                    except UnicodeDecodeError:
                        try:
                            # Fall back to latin-1 if UTF-8 fails
                            chunk_str = chunk.decode("latin-1")
                        except Exception as e:
                            logger.warning(f"Failed to decode chunk: {str(e)}")
                            continue

                    # Parse the chunk
                    soup.feed(chunk_str)

                    # Periodically clear the builder's memory
                    if len(soup.builder.text_chunks) > 1000:  # Arbitrary threshold
                        text_chunks.extend(soup.builder.text_chunks)
                        soup = BeautifulSoup("", parser, builder=StreamingHTMLParser())

            # Add any remaining text
            if hasattr(soup.builder, "text_chunks"):
                text_chunks.extend(soup.builder.text_chunks)

            # Join all text chunks with newlines
            content = "\n".join(text_chunks)
            logger.info(f"Extracted {len(content)} characters from large HTML file")

            return {
                "content": content,
                "type": "text",
                "source": "html",
                "processed_in_chunks": True,
                "original_size": file_path.stat().st_size,
                "processed_size": len(content),
            }

        except Exception as e:
            logger.error(f"Error processing large HTML file {file_path}: {str(e)}", exc_info=True)
            return {"error": f"Error processing large HTML file: {str(e)}"}

    def _process_large_text_file(self, file_path: Path, chunk_size: int = 1024 * 1024) -> Dict[str, Any]:
        """Process large text files in chunks to save memory."""
        try:
            logger.info(f"Processing large text file in chunks: {file_path}")
            content_parts = []

            with open(file_path, "r", encoding="utf-8") as f:
                while True:
                    chunk = f.read(chunk_size)
                    if not chunk:
                        break
                    content_parts.append(chunk)

            content = "".join(content_parts)
            logger.info(f"Read {len(content)} characters from large text file")
            return {"content": content, "type": "text", "processed_in_chunks": True}

        except Exception as e:
            logger.error(f"Error processing large text file {file_path}: {str(e)}", exc_info=True)
            return {"error": f"Error processing large text file: {str(e)}"}

        except UnicodeDecodeError:
            logger.warning(f"UTF-8 decode failed for {file_path}, trying latin-1")
            try:
                with open(file_path, "r", encoding="latin-1") as f:
                    content = f.read()
                    logger.info(f"Successfully read {len(content)} characters using latin-1 encoding")
                    return {"content": content, "type": "text", "encoding": "latin-1"}
            except Exception as e:
                error_msg = f"Failed to read text file: {str(e)}"
                logger.error(error_msg, exc_info=True)
                return {"error": error_msg}
        except Exception as e:
            error_msg = f"Error processing text file: {str(e)}"
            logger.error(error_msg, exc_info=True)
            return {"error": error_msg}

    def _process_image_file(self, file_path: Path) -> Dict[str, Any]:
        """Process an image file."""
        try:
            with Image.open(file_path) as img:
                return {
                    "type": "image",
                    "format": img.format,
                    "size": img.size,
                    "mode": img.mode,
                    "content": f"Image: {file_path.name} ({img.format}, {img.size[0]}x{img.size[1]}, {img.mode})",
                }
        except Exception as e:
            return {"error": f"Failed to process image: {str(e)}"}

    def _process_pdf_file(self, file_path: Path) -> Dict[str, Any]:
        """Process a PDF file."""
        try:
            text = []
            with fitz.open(file_path) as doc:
                for page in doc:
                    text.append(page.get_text())
            return {"content": "\n".join(text), "type": "document", "page_count": len(text)}
        except Exception as e:
            return {"error": f"Failed to process PDF: {str(e)}"}

    def _process_word_file(self, file_path: Path) -> Dict[str, Any]:
        """Process a Word document."""
        try:
            doc = docx.Document(file_path)
            text = [paragraph.text for paragraph in doc.paragraphs]
            return {"content": "\n".join(text), "type": "document", "paragraph_count": len(text)}
        except Exception as e:
            return {"error": f"Failed to process Word document: {str(e)}"}

    def _process_excel_file(self, file_path: Path) -> Dict[str, Any]:
        """Process an Excel file."""
        try:
            # Try pandas first
            try:
                df = pd.read_excel(file_path, sheet_name=None)
                content = []
                for sheet_name, sheet_data in df.items():
                    content.append(f"--- Sheet: {sheet_name} ---")
                    content.append(sheet_data.to_string())
                return {"content": "\n".join(content), "type": "spreadsheet", "sheet_count": len(df)}
            except:
                # Fallback to openpyxl
                wb = load_workbook(filename=file_path, read_only=True)
                content = []
                for sheet_name in wb.sheetnames:
                    content.append(f"--- Sheet: {sheet_name} ---")
                    sheet = wb[sheet_name]
                    for row in sheet.iter_rows(values_only=True):
                        content.append("\t".join(str(cell) if cell is not None else "" for cell in row))
                return {"content": "\n".join(content), "type": "spreadsheet", "sheet_count": len(wb.sheetnames)}
        except Exception as e:
            return {"error": f"Failed to process Excel file: {str(e)}"}

    def _process_powerpoint_file(self, file_path: Path) -> Dict[str, Any]:
        """Process a PowerPoint file."""
        try:
            prs = Presentation(file_path)
            content = []
            for i, slide in enumerate(prs.slides, 1):
                content.append(f"--- Slide {i} ---")
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        content.append(shape.text)
            return {"content": "\n".join(content), "type": "presentation", "slide_count": len(prs.slides)}
        except Exception as e:
            return {"error": f"Failed to process PowerPoint file: {str(e)}"}


class MemoryManager:
    """Manages memory operations including file processing and caching."""

    def __init__(self, cache_dir: Optional[Path] = None):
        self.cache_dir = cache_dir or Path(".cache/memory")
        self.cache_dir.mkdir(parents=True, exist_ok=True)
        self.processor = FileProcessor(self.cache_dir)

    def process_file(self, file_path: Union[str, Path]) -> Dict[str, Any]:
        """Process a file and return its content and metadata."""
        return self.processor.process_file(Path(file_path))

    def process_directory(
        self, dir_path: Union[str, Path], extensions: Optional[List[str]] = None
    ) -> List[Dict[str, Any]]:
        """Process all files in a directory matching the given extensions."""
        dir_path = Path(dir_path)
        if not dir_path.is_dir():
            return [{"error": f"Not a directory: {dir_path}"}]

        results = []
        for file_path in dir_path.rglob("*"):
            if file_path.is_file() and (not extensions or file_path.suffix.lower() in extensions):
                results.append(self.process_file(file_path))

        return results

    def clear_cache(self) -> None:
        """Clear the file processing cache."""
        import shutil

        if self.cache_dir.exists():
            shutil.rmtree(self.cache_dir)
        self.cache_dir.mkdir(parents=True)
